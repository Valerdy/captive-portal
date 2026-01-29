#!/usr/bin/env python3
"""Generate a visual PowerPoint presentation for the Captive Portal project.
Focus on diagrams, schemas and illustrations with minimal text."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Colors
DARK_BLUE = RGBColor(0x1A, 0x36, 0x5D)
MEDIUM_BLUE = RGBColor(0x1E, 0x40, 0xAF)
LIGHT_BLUE = RGBColor(0x25, 0x63, 0xEB)
ACCENT_BLUE = RGBColor(0xDB, 0xEA, 0xFE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x64, 0x74, 0x8B)
LIGHT_GRAY = RGBColor(0xF1, 0xF5, 0xF9)
GREEN = RGBColor(0x10, 0xB9, 0x81)
LIGHT_GREEN = RGBColor(0xD1, 0xFA, 0xE5)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
LIGHT_ORANGE = RGBColor(0xFE, 0xF3, 0xC7)
RED = RGBColor(0xEF, 0x44, 0x44)
LIGHT_RED = RGBColor(0xFE, 0xE2, 0xE2)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
LIGHT_PURPLE = RGBColor(0xED, 0xE9, 0xFE)
TEAL = RGBColor(0x14, 0xB8, 0xA6)
LIGHT_TEAL = RGBColor(0xCC, 0xFB, 0xF1)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


def add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, color, border_color=None, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(2)
    else:
        shape.line.fill.background()
    return shape


def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def set_text(shape, text, size=14, bold=False, color=BLACK, align=PP_ALIGN.LEFT):
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return tf


def add_para(tf, text, size=14, bold=False, color=BLACK, align=PP_ALIGN.LEFT, space_before=Pt(4)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    if space_before:
        p.space_before = space_before
    return p


def add_text(slide, left, top, width, height, text, size=14, bold=False, color=BLACK, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return tf


def header_bar(slide, title):
    add_rect(slide, 0, 0, W, Inches(1.1), MEDIUM_BLUE)
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12), Inches(0.8))
    set_text(txBox, title, size=28, bold=True, color=WHITE)


def footer_bar(slide):
    add_rect(slide, 0, Inches(7.44), W, Inches(0.06), LIGHT_BLUE)


def icon_box(slide, left, top, size, icon_text, bg_color, text_color=WHITE):
    """Draw a rounded square icon with an emoji/symbol."""
    shape = add_shape(slide, left, top, Inches(size), Inches(size), bg_color)
    set_text(shape, icon_text, size=int(size * 28), bold=True, color=text_color, align=PP_ALIGN.CENTER)
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    return shape


def arrow_right(slide, left, top, width=Inches(1), color=MEDIUM_BLUE):
    """Draw a right arrow."""
    arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, Inches(0.4))
    arr.fill.solid()
    arr.fill.fore_color.rgb = color
    arr.line.fill.background()
    return arr


def arrow_down(slide, left, top, height=Inches(0.6), color=MEDIUM_BLUE):
    """Draw a down arrow."""
    arr = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left, top, Inches(0.4), height)
    arr.fill.solid()
    arr.fill.fore_color.rgb = color
    arr.line.fill.background()
    return arr


def card(slide, left, top, width, height, title, subtitle, bg_color, border_color, icon_text="", title_color=DARK_BLUE):
    """Create a visual card with optional icon."""
    box = add_shape(slide, left, top, width, height, bg_color, border_color)
    y = top + Inches(0.15)
    if icon_text:
        add_text(slide, left, y, width, Inches(0.5), icon_text, size=24, align=PP_ALIGN.CENTER)
        y += Inches(0.5)
    add_text(slide, left + Inches(0.15), y, width - Inches(0.3), Inches(0.4),
             title, size=14, bold=True, color=title_color, align=PP_ALIGN.CENTER)
    if subtitle:
        add_text(slide, left + Inches(0.15), y + Inches(0.35), width - Inches(0.3), Inches(0.8),
                 subtitle, size=11, color=GRAY, align=PP_ALIGN.CENTER)
    return box


# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
add_rect(slide, 0, 0, W, Inches(0.08), LIGHT_BLUE)

# Big WiFi icon
wifi = add_shape(slide, Inches(5.4), Inches(1.2), Inches(2.5), Inches(2.5), MEDIUM_BLUE, shape_type=MSO_SHAPE.OVAL)
set_text(wifi, "📡", size=60, align=PP_ALIGN.CENTER)

add_text(slide, Inches(1), Inches(3.9), Inches(11.333), Inches(1),
         "Portail Captif WiFi — UCAC-ICAM", size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(5), Inches(11.333), Inches(0.6),
         "Projet de fin d'études — Janvier 2026", size=20, color=RGBColor(0xBF, 0xDB, 0xFE), align=PP_ALIGN.CENTER)
tf = add_text(slide, Inches(3), Inches(5.8), Inches(7.333), Inches(1),
              "Auteur : Valerdy", size=16, color=RGBColor(0xBF, 0xDB, 0xFE), align=PP_ALIGN.CENTER)
add_para(tf, "Encadrement technique avec Claude (Assistant IA)", size=14, color=GRAY, align=PP_ALIGN.CENTER)

add_rect(slide, 0, Inches(7.42), W, Inches(0.08), LIGHT_BLUE)

# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — Sommaire visuel (icons grid)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
header_bar(slide, "Sommaire")

sections = [
    ("🌐", "Architecture\nRéseau", LIGHT_BLUE, MEDIUM_BLUE),
    ("📡", "MikroTik\nHotspot", LIGHT_ORANGE, ORANGE),
    ("🔐", "FreeRADIUS\nAuthentification", LIGHT_GREEN, GREEN),
    ("🗄️", "PostgreSQL\nBase de données", LIGHT_PURPLE, PURPLE),
    ("⚡", "Flux\nComplet", LIGHT_TEAL, TEAL),
    ("🎯", "Conclusion", LIGHT_RED, RED),
]

for i, (icon, label, bg, border) in enumerate(sections):
    col = i % 3
    row = i // 3
    left = Inches(1.2 + col * 3.8)
    top = Inches(1.8 + row * 2.8)
    card(slide, left, top, Inches(3.2), Inches(2.2), label, "", bg, border, icon_text=icon)

footer_bar(slide)

# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — Architecture réseau (grand diagramme visuel)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
header_bar(slide, "Architecture Réseau Globale")

# Internet cloud at top center
inet = add_shape(slide, Inches(5.2), Inches(1.3), Inches(3), Inches(1), LIGHT_ORANGE, ORANGE, MSO_SHAPE.CLOUD)
add_text(slide, Inches(5.2), Inches(1.5), Inches(3), Inches(0.6),
         "☁  INTERNET", size=16, bold=True, color=BLACK, align=PP_ALIGN.CENTER)

# Arrow down from internet
arrow_down(slide, Inches(6.5), Inches(2.35), Inches(0.5), ORANGE)

# MikroTik router (center)
mk = add_shape(slide, Inches(4.2), Inches(2.9), Inches(5), Inches(1.5), ACCENT_BLUE, MEDIUM_BLUE)
tf = add_text(slide, Inches(4.2), Inches(3.0), Inches(5), Inches(0.5),
              "📡  MikroTik RB951Ui-2HnD", size=18, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
add_para(tf, "Hotspot  •  NAT  •  DHCP  •  DNS  •  Client RADIUS", size=12, color=GRAY, align=PP_ALIGN.CENTER)
add_para(tf, "10.242.18.6", size=13, bold=True, color=MEDIUM_BLUE, align=PP_ALIGN.CENTER)

# Arrow down-left to server
arrow_down(slide, Inches(4.5), Inches(4.45), Inches(0.5), GREEN)
# Arrow down-right to clients
arrow_down(slide, Inches(8.4), Inches(4.45), Inches(0.5), RED)

# Server (bottom left)
srv = add_shape(slide, Inches(0.5), Inches(5.1), Inches(5.5), Inches(2.0), LIGHT_GREEN, GREEN)
tf2 = add_text(slide, Inches(0.5), Inches(5.2), Inches(5.5), Inches(0.5),
               "🖥️  Serveur Backend", size=18, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
add_para(tf2, "Django REST  •  PostgreSQL  •  FreeRADIUS", size=12, color=GRAY, align=PP_ALIGN.CENTER)
add_para(tf2, "RADIUS UDP 1812/1813  |  API HTTP 8000", size=11, color=GRAY, align=PP_ALIGN.CENTER)
add_para(tf2, "Redis  •  Celery  •  Prometheus  •  Grafana", size=11, color=GRAY, align=PP_ALIGN.CENTER)

# Arrow between server and mikrotik (horizontal)
add_text(slide, Inches(2.2), Inches(4.55), Inches(2.5), Inches(0.4),
         "◄── RADIUS ──►", size=11, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

# Clients (bottom right)
cl = add_shape(slide, Inches(7.3), Inches(5.1), Inches(5.5), Inches(2.0), LIGHT_RED, RED)
tf3 = add_text(slide, Inches(7.3), Inches(5.2), Inches(5.5), Inches(0.5),
               "📱  Clients WiFi", size=18, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
add_para(tf3, "DHCP : 10.242.18.100 — 10.242.18.200", size=12, color=GRAY, align=PP_ALIGN.CENTER)
add_para(tf3, "Redirection captive → page de login", size=11, color=GRAY, align=PP_ALIGN.CENTER)
add_para(tf3, "Auth RADIUS (PAP/CHAP)", size=11, color=GRAY, align=PP_ALIGN.CENTER)

# Arrow between clients and mikrotik
add_text(slide, Inches(8.8), Inches(4.55), Inches(2.5), Inches(0.4),
         "◄── WiFi/DHCP ──►", size=11, bold=True, color=RED, align=PP_ALIGN.CENTER)

footer_bar(slide)

# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — Section MikroTik
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, 0, 0, W, Inches(0.06), LIGHT_BLUE)

circ = add_shape(slide, Inches(5.667), Inches(1.5), Inches(2), Inches(2), MEDIUM_BLUE, shape_type=MSO_SHAPE.OVAL)
set_text(circ, "📡", size=48, align=PP_ALIGN.CENTER)

add_text(slide, Inches(1), Inches(3.8), Inches(11.333), Inches(1),
         "Configuration MikroTik", size=36, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
add_text(slide, Inches(2), Inches(4.8), Inches(9.333), Inches(0.6),
         "Interfaces  •  DHCP  •  Hotspot  •  NAT  •  DNS  •  Sécurité", size=18, color=GRAY, align=PP_ALIGN.CENTER)

footer_bar(slide)

# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — MikroTik Interfaces (visual diagram)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
header_bar(slide, "MikroTik — Interfaces et Adressage IP")

# Central router box
router = add_shape(slide, Inches(4.5), Inches(2.5), Inches(4.3), Inches(2.5), ACCENT_BLUE, MEDIUM_BLUE)
add_text(slide, Inches(4.5), Inches(2.6), Inches(4.3), Inches(0.5),
         "📡 MikroTik RB951Ui-2HnD", size=16, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
add_text(slide, Inches(4.5), Inches(3.1), Inches(4.3), Inches(0.4),
         "RouterOS v6/v7", size=12, color=GRAY, align=PP_ALIGN.CENTER)

# ether1-WAN (left of router)
wan = add_shape(slide, Inches(0.5), Inches(1.5), Inches(3.2), Inches(1.3), LIGHT_ORANGE, ORANGE)
add_text(slide, Inches(0.5), Inches(1.6), Inches(3.2), Inches(0.4),
         "☁  ether1-WAN", size=15, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
add_text(slide, Inches(0.5), Inches(2.1), Inches(3.2), Inches(0.4),
         "Connexion Internet (WAN)", size=12, color=GRAY, align=PP_ALIGN.CENTER)

# Arrow WAN -> Router
arrow_right(slide, Inches(3.8), Inches(2.8), Inches(0.6), ORANGE)

# bridge-LAN box (inside router)
bridge = add_shape(slide, Inches(4.8), Inches(3.6), Inches(3.7), Inches(1.1), LIGHT_GREEN, GREEN)
add_text(slide, Inches(4.8), Inches(3.65), Inches(3.7), Inches(0.4),
         "🔗 bridge-LAN", size=13, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
add_text(slide, Inches(4.8), Inches(4.0), Inches(3.7), Inches(0.4),
         "10.242.18.6/24", size=12, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

# ether2-LAN (right)
lan = add_shape(slide, Inches(9.5), Inches(2.5), Inches(3.3), Inches(1.3), LIGHT_TEAL, TEAL)
add_text(slide, Inches(9.5), Inches(2.6), Inches(3.3), Inches(0.4),
         "🔌 ether2-LAN", size=15, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
add_text(slide, Inches(9.5), Inches(3.1), Inches(3.3), Inches(0.4),
         "Réseau local (LAN)", size=12, color=GRAY, align=PP_ALIGN.CENTER)

# Arrow Router -> LAN
arrow_right(slide, Inches(8.85), Inches(2.8), Inches(0.6), TEAL)

# DHCP Pool info (bottom)
pool = add_shape(slide, Inches(3.5), Inches(5.5), Inches(6.3), Inches(1.2), LIGHT_PURPLE, PURPLE)
add_text(slide, Inches(3.5), Inches(5.6), Inches(6.3), Inches(0.4),
         "🏊 Pool DHCP : 10.242.18.100 — 10.242.18.200", size=14, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
add_text(slide, Inches(3.5), Inches(6.05), Inches(6.3), Inches(0.4),
         "101 adresses  •  Gateway: 10.242.18.6  •  DNS: 10.242.18.6", size=12, color=GRAY, align=PP_ALIGN.CENTER)

# Arrow bridge -> pool
arrow_down(slide, Inches(6.5), Inches(4.8), Inches(0.6), PURPLE)

footer_bar(slide)

# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — DHCP (visual flow)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
header_bar(slide, "MikroTik — Serveur DHCP")

# DHCP flow: 4 steps horizontally
steps = [
    ("1", "📱 Client\nse connecte", "Le client WiFi\nrejoint le réseau", LIGHT_BLUE, MEDIUM_BLUE),
    ("2", "📨 DHCP\nDiscover", "Demande d'adresse\nIP (broadcast)", LIGHT_ORANGE, ORANGE),
    ("3", "📡 DHCP\nOffer", "Le routeur propose\nune IP du pool", LIGHT_GREEN, GREEN),
    ("4", "✅ IP\nAttribuée", "10.242.18.x\nGateway + DNS", LIGHT_PURPLE, PURPLE),
]

for i, (num, title, desc, bg, border) in enumerate(steps):
    left = Inches(0.5 + i * 3.2)
    top = Inches(1.6)
    # Number circle
    circ = add_shape(slide, left + Inches(1), top, Inches(0.7), Inches(0.7), border, shape_type=MSO_SHAPE.OVAL)
    set_text(circ, num, size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Card
    box = add_shape(slide, left, top + Inches(0.9), Inches(2.7), Inches(2.0), bg, border)
    add_text(slide, left + Inches(0.1), top + Inches(1.0), Inches(2.5), Inches(0.6),
             title, size=15, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
    add_text(slide, left + Inches(0.1), top + Inches(1.8), Inches(2.5), Inches(0.8),
             desc, size=12, color=GRAY, align=PP_ALIGN.CENTER)
    # Arrow between steps
    if i < 3:
        arrow_right(slide, left + Inches(2.75), top + Inches(1.6), Inches(0.4), border)

# Config summary at bottom
cfg = add_shape(slide, Inches(1), Inches(5.5), Inches(11.333), Inches(1.3), LIGHT_GRAY, MEDIUM_BLUE)
add_text(slide, Inches(1.3), Inches(5.6), Inches(10.7), Inches(0.4),
         "Configuration DHCP", size=16, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
tf = add_text(slide, Inches(1.3), Inches(6.05), Inches(10.7), Inches(0.6),
              "Pool : dhcp-pool  •  Plage : .100 à .200  •  Interface : bridge-LAN  •  Réseau : 10.242.18.0/24",
              size=13, color=GRAY, align=PP_ALIGN.CENTER)

footer_bar(slide)

# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — Hotspot (visual flow)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
header_bar(slide, "MikroTik — Principe du Hotspot")

# Visual: Client -> Redirect -> Login -> RADIUS -> Internet
steps_h = [
    ("📱", "Client WiFi", "Connexion au\nréseau UCAC", LIGHT_BLUE, MEDIUM_BLUE),
    ("🔄", "Redirection", "Trafic HTTP\nintercepté", LIGHT_ORANGE, ORANGE),
    ("🔑", "Page Login", "wifi.ucac-icam.cm\n/login", LIGHT_PURPLE, PURPLE),
    ("🔐", "RADIUS", "Vérification\ncredentials", LIGHT_GREEN, GREEN),
    ("🌐", "Internet", "Accès autorisé\navec QoS", LIGHT_TEAL, TEAL),
]

for i, (icon, title, desc, bg, border) in enumerate(steps_h):
    left = Inches(0.3 + i * 2.6)
    top = Inches(1.5)
    # Icon circle
    circ = add_shape(slide, left + Inches(0.6), top, Inches(1.2), Inches(1.2), border, shape_type=MSO_SHAPE.OVAL)
    set_text(circ, icon, size=32, align=PP_ALIGN.CENTER)
    # Label
    add_text(slide, left, top + Inches(1.3), Inches(2.4), Inches(0.5),
             title, size=14, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
    add_text(slide, left, top + Inches(1.8), Inches(2.4), Inches(0.6),
             desc, size=11, color=GRAY, align=PP_ALIGN.CENTER)
    # Arrow
    if i < 4:
        arrow_right(slide, left + Inches(2.15), top + Inches(0.4), Inches(0.4), border)

# Hotspot profile details
prof = add_shape(slide, Inches(0.5), Inches(4.2), Inches(6), Inches(2.8), ACCENT_BLUE, MEDIUM_BLUE)
add_text(slide, Inches(0.7), Inches(4.3), Inches(5.6), Inches(0.5),
         "📋 Profil : ucac-icam-profile", size=16, bold=True, color=DARK_BLUE, align=PP_ALIGN.LEFT)
tf = add_text(slide, Inches(0.7), Inches(4.85), Inches(5.6), Inches(2),
              "• Adresse : 10.242.18.6", size=13, color=BLACK, align=PP_ALIGN.LEFT)
add_para(tf, "• DNS : wifi.ucac-icam.cm", size=13, color=BLACK)
add_para(tf, "• RADIUS : activé (auth + accounting)", size=13, color=BLACK)
add_para(tf, "• Interim update : 5 minutes", size=13, color=BLACK)
add_para(tf, "• Login : HTTP-CHAP + HTTP-PAP", size=13, color=BLACK)

# Walled garden
wg = add_shape(slide, Inches(7), Inches(4.2), Inches(5.8), Inches(2.8), LIGHT_GREEN, GREEN)
add_text(slide, Inches(7.2), Inches(4.3), Inches(5.4), Inches(0.5),
         "🏰 Walled Garden", size=16, bold=True, color=DARK_BLUE, align=PP_ALIGN.LEFT)
tf2 = add_text(slide, Inches(7.2), Inches(4.85), Inches(5.4), Inches(2),
               "Pages accessibles SANS authentification :", size=13, color=BLACK, align=PP_ALIGN.LEFT)
add_para(tf2, "• Serveur backend Django", size=13, color=BLACK)
add_para(tf2, "• API d'enregistrement", size=13, color=BLACK)
add_para(tf2, "• Page de login personnalisée", size=13, color=BLACK)
add_para(tf2, "", size=8, color=BLACK)
add_para(tf2, "Le hotspot crée auto. les règles\nfirewall, NAT et queues", size=12, color=GRAY)

footer_bar(slide)

# ═══════════════════════════════════════════════════════════════
# SLIDE 8 — NAT Masquerade (visual diagram)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
header_bar(slide, "MikroTik — NAT Masquerade")

# Visual: Private IP -> Router (NAT) -> Public IP -> Internet
# Client
cl_box = add_shape(slide, Inches(0.5), Inches(2.5), Inches(3), Inches(2.5), LIGHT_BLUE, MEDIUM_BLUE)
add_text(slide, Inches(0.5), Inches(2.6), Inches(3), Inches(0.5),
         "📱 Client WiFi", size=16, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
add_text(slide, Inches(0.5), Inches(3.2), Inches(3), Inches(0.4),
         "IP privée", size=13, color=GRAY, align=PP_ALIGN.CENTER)
add_text(slide, Inches(0.5), Inches(3.6), Inches(3), Inches(0.5),
         "10.242.18.x", size=18, bold=True, color=MEDIUM_BLUE, align=PP_ALIGN.CENTER)
add_text(slide, Inches(0.5), Inches(4.2), Inches(3), Inches(0.4),
         "Source: 10.242.18.105:4532", size=11, color=GRAY, align=PP_ALIGN.CENTER)

# Arrow
arrow_right(slide, Inches(3.6), Inches(3.5), Inches(0.8), MEDIUM_BLUE)

# Router NAT
nat_box = add_shape(slide, Inches(4.5), Inches(2.2), Inches(4.3), Inches(3.2), LIGHT_ORANGE, ORANGE)
add_text(slide, Inches(4.5), Inches(2.3), Inches(4.3), Inches(0.5),
         "🔄 MikroTik — NAT", size=16, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
add_text(slide, Inches(4.5), Inches(2.9), Inches(4.3), Inches(0.4),
         "Masquerade", size=22, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

# Visual: IP transformation
tf_nat = add_text(slide, Inches(4.7), Inches(3.5), Inches(3.9), Inches(0.4),
                  "10.242.18.105  →  IP publique WAN", size=13, color=BLACK, align=PP_ALIGN.CENTER)
add_para(tf_nat, "", size=6, color=BLACK)
add_para(tf_nat, "srcnat • out-interface=ether1-WAN", size=11, color=GRAY, align=PP_ALIGN.CENTER)
add_para(tf_nat, "action=masquerade", size=11, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

# Arrow
arrow_right(slide, Inches(8.9), Inches(3.5), Inches(0.8), GREEN)

# Internet
inet_box = add_shape(slide, Inches(9.8), Inches(2.5), Inches(3), Inches(2.5), LIGHT_GREEN, GREEN)
add_text(slide, Inches(9.8), Inches(2.6), Inches(3), Inches(0.5),
         "☁ Internet", size=16, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
add_text(slide, Inches(9.8), Inches(3.2), Inches(3), Inches(0.4),
         "IP publique", size=13, color=GRAY, align=PP_ALIGN.CENTER)
add_text(slide, Inches(9.8), Inches(3.6), Inches(3), Inches(0.5),
         "WAN IP", size=18, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
add_text(slide, Inches(9.8), Inches(4.2), Inches(3), Inches(0.4),
         "Source: IP_WAN:random", size=11, color=GRAY, align=PP_ALIGN.CENTER)

# Explanation
expl = add_shape(slide, Inches(1.5), Inches(5.8), Inches(10.3), Inches(1.0), LIGHT_GRAY, MEDIUM_BLUE)
add_text(slide, Inches(1.7), Inches(5.9), Inches(9.9), Inches(0.7),
         "Sans NAT, les clients en 10.242.18.x ne peuvent pas accéder à Internet. Le masquerade remplace l'IP source privée par l'IP publique du routeur.",
         size=13, color=GRAY, align=PP_ALIGN.CENTER)

footer_bar(slide)

# ═══════════════════════════════════════════════════════════════
# SLIDE 9 — DNS & Sécurité (visual)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
header_bar(slide, "MikroTik — DNS et Sécurité")

# Left: DNS Configuration
dns_box = add_shape(slide, Inches(0.5), Inches(1.4), Inches(6), Inches(2.8), LIGHT_BLUE, MEDIUM_BLUE)
add_text(slide, Inches(0.7), Inches(1.5), Inches(5.6), Inches(0.5),
         "🌐 Serveur DNS Local", size=16, bold=True, color=DARK_BLUE, align=PP_ALIGN.LEFT)
tf = add_text(slide, Inches(0.7), Inches(2.05), Inches(5.6), Inches(2),
              "• Résolution DNS sur le routeur", size=13, color=BLACK)
add_para(tf, "• Upstream : 8.8.8.8 + 1.1.1.1", size=13, color=BLACK)
add_para(tf, "• Cache : 4096 KiB", size=13, color=BLACK)
add_para(tf, "• Toutes les requêtes DNS redirigées localement", size=13, color=BLACK)

# Right: Blocking
block_box = add_shape(slide, Inches(6.9), Inches(1.4), Inches(5.9), Inches(2.8), LIGHT_RED, RED)
add_text(slide, Inches(7.1), Inches(1.5), Inches(5.5), Inches(0.5),
         "🚫 Blocage de Sites", size=16, bold=True, color=DARK_BLUE, align=PP_ALIGN.LEFT)
tf2 = add_text(slide, Inches(7.1), Inches(2.05), Inches(5.5), Inches(2),
               "• DNS statique → 0.0.0.0", size=13, color=BLACK)
add_para(tf2, "• Support regex (wildcard)", size=13, color=BLACK)
add_para(tf2, "• Géré depuis le backend Django", size=13, color=BLACK)
add_para(tf2, "• Ex: facebook.com → 0.0.0.0", size=13, color=BLACK)

# Bottom: DNS forcing diagram
force_box = add_shape(slide, Inches(0.5), Inches(4.5), Inches(6), Inches(2.5), LIGHT_ORANGE, ORANGE)
add_text(slide, Inches(0.7), Inches(4.6), Inches(5.6), Inches(0.5),
         "🔒 Forçage DNS", size=16, bold=True, color=DARK_BLUE, align=PP_ALIGN.LEFT)
tf3 = add_text(slide, Inches(0.7), Inches(5.15), Inches(5.6), Inches(1.5),
               "• Redirection port 53 (UDP+TCP) → routeur", size=13, color=BLACK)
add_para(tf3, "• Blocage DNS over TLS (port 853)", size=13, color=BLACK)
add_para(tf3, "• Empêche le contournement du filtrage", size=13, color=BLACK)

# API RouterOS
api_box = add_shape(slide, Inches(6.9), Inches(4.5), Inches(5.9), Inches(2.5), LIGHT_PURPLE, PURPLE)
add_text(slide, Inches(7.1), Inches(4.6), Inches(5.5), Inches(0.5),
         "⚙️ API RouterOS", size=16, bold=True, color=DARK_BLUE, align=PP_ALIGN.LEFT)
tf4 = add_text(slide, Inches(7.1), Inches(5.15), Inches(5.5), Inches(1.5),
               "• Port 8728 (API)", size=13, color=BLACK)
add_para(tf4, "• Agent Node.js dédié", size=13, color=BLACK)
add_para(tf4, "• Communication backend ↔ routeur", size=13, color=BLACK)

footer_bar(slide)

# ═══════════════════════════════════════════════════════════════
# SLIDE 10 — Section FreeRADIUS
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, 0, 0, W, Inches(0.06), LIGHT_BLUE)

circ = add_shape(slide, Inches(5.667), Inches(1.5), Inches(2), Inches(2), GREEN, shape_type=MSO_SHAPE.OVAL)
set_text(circ, "🔐", size=48, align=PP_ALIGN.CENTER)

add_text(slide, Inches(1), Inches(3.8), Inches(11.333), Inches(1),
         "Configuration FreeRADIUS", size=36, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
add_text(slide, Inches(2), Inches(4.8), Inches(9.333), Inches(0.6),
         "Serveur AAA  •  Authentication  •  Authorization  •  Accounting", size=18, color=GRAY, align=PP_ALIGN.CENTER)

footer_bar(slide)

# ═══════════════════════════════════════════════════════════════
# SLIDE 11 — FreeRADIUS AAA (3 colonnes visuelles)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
header_bar(slide, "FreeRADIUS — Authentification, Autorisation, Accounting")

# 3 big cards for A, A, A
aaa = [
    ("🔑", "Authentication", "Vérifier l'identité",
     ["Username + Password", "Table radcheck (PostgreSQL)", "Statut actif vérifié", "PAP / CHAP supportés"],
     LIGHT_GREEN, GREEN),
    ("📋", "Authorization", "Définir les droits",
     ["Bande passante (Rate-Limit)", "Durée session (Timeout)", "Tables radreply + radgroupreply", "Ex: 5M/10M, 8h session"],
     LIGHT_BLUE, MEDIUM_BLUE),
    ("📊", "Accounting", "Tracer l'activité",
     ["Sessions : début, durée, fin", "Octets transférés (up/down)", "Table radacct", "Mise à jour toutes les 5 min"],
     LIGHT_PURPLE, PURPLE),
]

for i, (icon, title, subtitle, items, bg, border) in enumerate(aaa):
    left = Inches(0.4 + i * 4.3)
    top = Inches(1.4)
    box = add_shape(slide, left, top, Inches(3.8), Inches(5.5), bg, border)
    # Icon
    circ = add_shape(slide, left + Inches(1.3), top + Inches(0.2), Inches(1.2), Inches(1.2), border, shape_type=MSO_SHAPE.OVAL)
    set_text(circ, icon, size=30, align=PP_ALIGN.CENTER)
    # Title
    add_text(slide, left + Inches(0.2), top + Inches(1.5), Inches(3.4), Inches(0.5),
             title, size=18, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
    add_text(slide, left + Inches(0.2), top + Inches(2.0), Inches(3.4), Inches(0.4),
             subtitle, size=13, color=GRAY, align=PP_ALIGN.CENTER)
    # Items
    y = top + Inches(2.6)
    for item in items:
        add_text(slide, left + Inches(0.3), y, Inches(3.2), Inches(0.35),
                 "•  " + item, size=12, color=BLACK, align=PP_ALIGN.LEFT)
        y += Inches(0.35)

footer_bar(slide)

# ═══════════════════════════════════════════════════════════════
# SLIDE 12 — FreeRADIUS flux auth (visual flow)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
header_bar(slide, "FreeRADIUS — Flux d'Authentification")

# Vertical flow
flow_steps = [
    ("1", "📱 Client envoie credentials", "Page de login hotspot (HTTP-CHAP / HTTP-PAP)", LIGHT_BLUE, MEDIUM_BLUE),
    ("2", "📡 MikroTik → Access-Request", "User-Name, Password, NAS-IP, MAC  •  UDP 1812", LIGHT_ORANGE, ORANGE),
    ("3", "🔍 FreeRADIUS → authorize (SQL)", "SELECT FROM radcheck WHERE username=? AND statut=true", LIGHT_GREEN, GREEN),
    ("4", "🔑 FreeRADIUS → authenticate", "Comparaison mot de passe (PAP: clair, CHAP: MD5)", LIGHT_PURPLE, PURPLE),
]

for i, (num, title, desc, bg, border) in enumerate(flow_steps):
    left_num = Inches(0.5)
    left_card = Inches(1.5)
    top = Inches(1.4 + i * 1.4)
    # Number circle
    circ = add_shape(slide, left_num, top + Inches(0.1), Inches(0.7), Inches(0.7), border, shape_type=MSO_SHAPE.OVAL)
    set_text(circ, num, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Card
    box = add_shape(slide, left_card, top, Inches(7), Inches(1.0), bg, border)
    add_text(slide, left_card + Inches(0.2), top + Inches(0.05), Inches(6.6), Inches(0.4),
             title, size=14, bold=True, color=DARK_BLUE)
    add_text(slide, left_card + Inches(0.2), top + Inches(0.5), Inches(6.6), Inches(0.4),
             desc, size=11, color=GRAY)

# Result boxes on the right
accept = add_shape(slide, Inches(9), Inches(1.5), Inches(3.8), Inches(2.2), LIGHT_GREEN, GREEN)
add_text(slide, Inches(9), Inches(1.6), Inches(3.8), Inches(0.5),
         "✅ Access-Accept", size=16, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
tf_a = add_text(slide, Inches(9.2), Inches(2.2), Inches(3.4), Inches(1.3),
                "Attributs retournés :", size=12, bold=True, color=DARK_BLUE)
add_para(tf_a, "• Mikrotik-Rate-Limit = 5M/10M", size=11, color=BLACK)
add_para(tf_a, "• Session-Timeout = 28800", size=11, color=BLACK)
add_para(tf_a, "• Idle-Timeout = 600", size=11, color=BLACK)

reject = add_shape(slide, Inches(9), Inches(4.2), Inches(3.8), Inches(1.5), LIGHT_RED, RED)
add_text(slide, Inches(9), Inches(4.3), Inches(3.8), Inches(0.5),
         "❌ Access-Reject", size=16, bold=True, color=RED, align=PP_ALIGN.CENTER)
tf_r = add_text(slide, Inches(9.2), Inches(4.9), Inches(3.4), Inches(0.6),
                "Log dans radpostauth", size=12, color=BLACK)
add_para(tf_r, "reply = 'Access-Reject'", size=11, color=GRAY)

footer_bar(slide)

# ═══════════════════════════════════════════════════════════════
# SLIDE 13 — FreeRADIUS config files (visual cards)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
header_bar(slide, "FreeRADIUS — Fichiers de Configuration")

files = [
    ("📄", "clients.conf", "Client RADIUS (MikroTik)", "IP: 10.242.18.6\nSecret partagé\nnastype: other", LIGHT_ORANGE, ORANGE),
    ("🗃️", "mods-enabled/sql", "Driver PostgreSQL", "rlm_sql_postgresql\nlocalhost:5432\ncaptive_portal", LIGHT_GREEN, GREEN),
    ("⚙️", "sites-available/default", "Sites RADIUS", "Port 1812 (auth)\nPort 1813 (acct)\nPAP + CHAP", LIGHT_BLUE, MEDIUM_BLUE),
    ("📖", "dictionary.mikrotik", "Attributs MikroTik", "Vendor ID: 14988\nRate-Limit (ID 8)\nTotal-Limit...", LIGHT_PURPLE, PURPLE),
]

for i, (icon, name, subtitle, desc, bg, border) in enumerate(files):
    col = i % 2
    row = i // 2
    left = Inches(0.5 + col * 6.4)
    top = Inches(1.4 + row * 2.9)
    box = add_shape(slide, left, top, Inches(5.9), Inches(2.5), bg, border)
    # Icon
    circ = add_shape(slide, left + Inches(0.3), top + Inches(0.3), Inches(0.9), Inches(0.9), border, shape_type=MSO_SHAPE.OVAL)
    set_text(circ, icon, size=22, align=PP_ALIGN.CENTER)
    # Title
    add_text(slide, left + Inches(1.5), top + Inches(0.3), Inches(4.1), Inches(0.4),
             name, size=16, bold=True, color=DARK_BLUE)
    add_text(slide, left + Inches(1.5), top + Inches(0.7), Inches(4.1), Inches(0.4),
             subtitle, size=12, color=GRAY)
    # Description
    add_text(slide, left + Inches(1.5), top + Inches(1.2), Inches(4.1), Inches(1.2),
             desc, size=12, color=BLACK)

footer_bar(slide)

# ═══════════════════════════════════════════════════════════════
# SLIDE 14 — Section PostgreSQL
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, 0, 0, W, Inches(0.06), LIGHT_BLUE)

circ = add_shape(slide, Inches(5.667), Inches(1.5), Inches(2), Inches(2), PURPLE, shape_type=MSO_SHAPE.OVAL)
set_text(circ, "🗄️", size=48, align=PP_ALIGN.CENTER)

add_text(slide, Inches(1), Inches(3.8), Inches(11.333), Inches(1),
         "Configuration PostgreSQL", size=36, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
add_text(slide, Inches(2), Inches(4.8), Inches(9.333), Inches(0.6),
         "Base de données captive_portal  •  Tables RADIUS  •  Tables Django", size=18, color=GRAY, align=PP_ALIGN.CENTER)

footer_bar(slide)

# ═══════════════════════════════════════════════════════════════
# SLIDE 15 — PostgreSQL tables (visual schema)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
header_bar(slide, "PostgreSQL — Schéma des Tables")

# Left: RADIUS tables
add_text(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(0.5),
         "🔐 Tables RADIUS (FreeRADIUS)", size=16, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

radius_tables = [
    ("radcheck", "Auth: username, password, statut"),
    ("radreply", "Réponse: attributs individuels"),
    ("radusergroup", "Liens utilisateur → groupe"),
    ("radgroupreply", "Réponse par groupe (QoS)"),
    ("radacct", "Sessions: durée, octets, MAC"),
    ("radpostauth", "Journal des authentifications"),
]

for i, (name, desc) in enumerate(radius_tables):
    top = Inches(1.85 + i * 0.82)
    box = add_shape(slide, Inches(0.5), top, Inches(6), Inches(0.7), LIGHT_GREEN, GREEN)
    add_text(slide, Inches(0.7), top + Inches(0.05), Inches(2.2), Inches(0.35),
             name, size=13, bold=True, color=DARK_BLUE)
    add_text(slide, Inches(3), top + Inches(0.05), Inches(3.3), Inches(0.35),
             desc, size=11, color=GRAY)

# Right: Django tables
add_text(slide, Inches(6.9), Inches(1.3), Inches(6), Inches(0.5),
         "🖥️ Tables Django (Gestion)", size=16, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

django_tables = [
    ("core_user", "Utilisateurs (extends AbstractUser)"),
    ("core_profile", "Profils (quotas, bande passante)"),
    ("core_promotion", "Groupes (classes, départements)"),
    ("core_userquota", "Quotas utilisateurs"),
    ("core_blockedsite", "Sites bloqués (DNS)"),
    ("mikrotik_config", "Config routeur MikroTik"),
]

for i, (name, desc) in enumerate(django_tables):
    top = Inches(1.85 + i * 0.82)
    box = add_shape(slide, Inches(6.9), top, Inches(6), Inches(0.7), LIGHT_PURPLE, PURPLE)
    add_text(slide, Inches(7.1), top + Inches(0.05), Inches(2.2), Inches(0.35),
             name, size=13, bold=True, color=DARK_BLUE)
    add_text(slide, Inches(9.4), top + Inches(0.05), Inches(3.3), Inches(0.35),
             desc, size=11, color=GRAY)

# Connection arrow
add_text(slide, Inches(4.5), Inches(7.0), Inches(4.3), Inches(0.4),
         "◄── Synchronisation Django ↔ RADIUS ──►", size=12, bold=True, color=MEDIUM_BLUE, align=PP_ALIGN.CENTER)

footer_bar(slide)

# ═══════════════════════════════════════════════════════════════
# SLIDE 16 — PostgreSQL ↔ FreeRADIUS link (visual)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
header_bar(slide, "PostgreSQL — Lien FreeRADIUS ↔ Django")

# 3 phases as visual flow
phases = [
    ("🔍", "Authorize", "Vérification", "SELECT FROM radcheck\nWHERE username=?\nAND statut=true",
     LIGHT_GREEN, GREEN),
    ("📝", "Post-Auth", "Journalisation", "INSERT INTO radpostauth\n(username, reply, authdate)\nSuccès ou Rejet",
     LIGHT_ORANGE, ORANGE),
    ("📊", "Accounting", "Comptabilité", "Start → INSERT radacct\nInterim → UPDATE octets\nStop → UPDATE fin session",
     LIGHT_BLUE, MEDIUM_BLUE),
]

for i, (icon, title, subtitle, desc, bg, border) in enumerate(phases):
    left = Inches(0.4 + i * 4.3)
    top = Inches(1.4)
    box = add_shape(slide, left, top, Inches(3.8), Inches(3.5), bg, border)
    circ = add_shape(slide, left + Inches(1.3), top + Inches(0.2), Inches(1.2), Inches(1.2), border, shape_type=MSO_SHAPE.OVAL)
    set_text(circ, icon, size=28, align=PP_ALIGN.CENTER)
    add_text(slide, left + Inches(0.2), top + Inches(1.5), Inches(3.4), Inches(0.5),
             title, size=16, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
    add_text(slide, left + Inches(0.2), top + Inches(1.95), Inches(3.4), Inches(0.35),
             subtitle, size=12, color=GRAY, align=PP_ALIGN.CENTER)
    add_text(slide, left + Inches(0.3), top + Inches(2.4), Inches(3.2), Inches(1),
             desc, size=12, color=BLACK, align=PP_ALIGN.CENTER)

# Django sync box
sync = add_shape(slide, Inches(1), Inches(5.3), Inches(11.333), Inches(1.7), LIGHT_PURPLE, PURPLE)
add_text(slide, Inches(1.2), Inches(5.4), Inches(10.9), Inches(0.5),
         "🔄 Synchronisation Django → RADIUS", size=16, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
tf = add_text(slide, Inches(1.5), Inches(5.95), Inches(10.3), Inches(0.8),
              "Admin active un utilisateur → Django crée radcheck + radreply + radusergroup", size=13, color=BLACK, align=PP_ALIGN.CENTER)
add_para(tf, "Profil modifié → Django met à jour radgroupreply + radgroupcheck  •  Retry automatique en cas d'échec",
         size=12, color=GRAY, align=PP_ALIGN.CENTER)

footer_bar(slide)

# ═══════════════════════════════════════════════════════════════
# SLIDE 17 — Flux complet (grand diagramme)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
header_bar(slide, "Flux d'Authentification Complet")

# 8 steps in a visual flow (2 rows of 4)
steps_flow = [
    ("1", "📱", "Connexion WiFi", "IP via DHCP\n10.242.18.x", LIGHT_BLUE, MEDIUM_BLUE),
    ("2", "🔄", "Redirection", "HTTP intercepté\n→ page login", LIGHT_ORANGE, ORANGE),
    ("3", "🔑", "Saisie credentials", "Username +\nPassword", LIGHT_PURPLE, PURPLE),
    ("4", "📡", "Access-Request", "MikroTik → RADIUS\nUDP 1812", LIGHT_TEAL, TEAL),
    ("5", "🗄️", "PostgreSQL", "radcheck :\nstatut + password", LIGHT_PURPLE, PURPLE),
    ("6", "✅", "Access-Accept", "Rate-Limit\nTimeout", LIGHT_GREEN, GREEN),
    ("7", "🌐", "Accès Internet", "QoS appliquée\n5M/10M", LIGHT_BLUE, MEDIUM_BLUE),
    ("8", "📊", "Accounting", "Sessions radacct\nToutes les 5 min", LIGHT_ORANGE, ORANGE),
]

for i, (num, icon, title, desc, bg, border) in enumerate(steps_flow):
    row = i // 4
    col = i % 4
    left = Inches(0.3 + col * 3.25)
    top = Inches(1.4 + row * 3.0)
    # Card
    box = add_shape(slide, left, top, Inches(2.8), Inches(2.5), bg, border)
    # Number
    circ = add_shape(slide, left + Inches(0.1), top + Inches(0.1), Inches(0.5), Inches(0.5), border, shape_type=MSO_SHAPE.OVAL)
    set_text(circ, num, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Icon
    add_text(slide, left, top + Inches(0.15), Inches(2.8), Inches(0.5),
             icon, size=28, align=PP_ALIGN.CENTER)
    # Title
    add_text(slide, left + Inches(0.1), top + Inches(0.8), Inches(2.6), Inches(0.5),
             title, size=14, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
    # Desc
    add_text(slide, left + Inches(0.1), top + Inches(1.35), Inches(2.6), Inches(0.9),
             desc, size=11, color=GRAY, align=PP_ALIGN.CENTER)
    # Arrow
    if col < 3:
        arrow_right(slide, left + Inches(2.85), top + Inches(1.0), Inches(0.35), border)

# Arrow from row 1 to row 2 (down)
arrow_down(slide, Inches(12.25), Inches(3.95), Inches(0.4), TEAL)

footer_bar(slide)

# ═══════════════════════════════════════════════════════════════
# SLIDE 18 — Conclusion (visual summary)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
header_bar(slide, "Conclusion")

components = [
    ("📡", "MikroTik", "Hotspot • NAT • DHCP\nDNS • Client RADIUS", LIGHT_ORANGE, ORANGE),
    ("🔐", "FreeRADIUS", "Serveur AAA central\nPAP/CHAP via SQL", LIGHT_GREEN, GREEN),
    ("🗄️", "PostgreSQL", "Tables RADIUS + Django\nUtilisateurs & Sessions", LIGHT_PURPLE, PURPLE),
    ("⚙️", "Django REST", "API de gestion\nSync RADIUS & MikroTik", LIGHT_BLUE, MEDIUM_BLUE),
    ("🖥️", "Vue.js 3", "Dashboard admin\nTemps réel", LIGHT_TEAL, TEAL),
]

for i, (icon, title, desc, bg, border) in enumerate(components):
    left = Inches(0.3 + i * 2.6)
    top = Inches(1.5)
    box = add_shape(slide, left, top, Inches(2.3), Inches(3.0), bg, border)
    circ = add_shape(slide, left + Inches(0.55), top + Inches(0.2), Inches(1.2), Inches(1.2), border, shape_type=MSO_SHAPE.OVAL)
    set_text(circ, icon, size=30, align=PP_ALIGN.CENTER)
    add_text(slide, left + Inches(0.1), top + Inches(1.5), Inches(2.1), Inches(0.4),
             title, size=14, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
    add_text(slide, left + Inches(0.1), top + Inches(1.95), Inches(2.1), Inches(0.8),
             desc, size=11, color=GRAY, align=PP_ALIGN.CENTER)

# Result box
result = add_shape(slide, Inches(1.5), Inches(5.0), Inches(10.333), Inches(1.8), ACCENT_BLUE, MEDIUM_BLUE)
add_text(slide, Inches(1.7), Inches(5.1), Inches(9.9), Inches(0.5),
         "🎯 Résultat", size=20, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1.7), Inches(5.65), Inches(9.9), Inches(0.9),
         "Un système WiFi sécurisé, évolutif et adapté au contexte universitaire UCAC-ICAM\nOrchestration Docker Compose de 11 services",
         size=15, color=GRAY, align=PP_ALIGN.CENTER)

footer_bar(slide)

# ═══════════════════════════════════════════════════════════════
# SLIDE 19 — Merci
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
add_rect(slide, 0, 0, W, Inches(0.08), LIGHT_BLUE)

add_text(slide, Inches(1), Inches(2.5), Inches(11.333), Inches(1.2),
         "Merci pour votre attention", size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(3.8), Inches(11.333), Inches(0.8),
         "Questions ?", size=28, color=RGBColor(0xBF, 0xDB, 0xFE), align=PP_ALIGN.CENTER)

add_text(slide, Inches(3), Inches(5.5), Inches(7.333), Inches(0.5),
         "Portail Captif WiFi — UCAC-ICAM — Janvier 2026", size=14, color=GRAY, align=PP_ALIGN.CENTER)

add_rect(slide, 0, Inches(7.42), W, Inches(0.08), LIGHT_BLUE)

# Save
output_path = "/home/user/captive-portal/Presentation_Portail_Captif_UCAC_ICAM.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
